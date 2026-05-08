"""
Renderização do mural em ``.pptx`` editável via python-pptx.

Recebe o mesmo :class:`LayoutMural` usado pelo PNG e produz um PowerPoint
no qual:

- O slide tem o tamanho exato do template (em EMU).
- O template entra como imagem de fundo, ocupando o slide inteiro.
- Cada texto (mês, dia, nome, cargo) vira UMA caixa de texto independente,
  posicionada na mesma coordenada usada pelo Pillow — assim o PPTX bate
  visualmente com o PNG.
- Se ``cfg.PPTX_GRUPO_POR_COLUNA`` for ``True``, todas as caixas de uma
  mesma coluna ficam num grupo (``GroupShape``) — o que permite arrastar a
  coluna inteira no Canva ou no PowerPoint.

python-pptx (resumo prático):
    - ``Presentation()`` cria uma apresentação vazia.
    - ``prs.slide_width / slide_height`` definem o slide em EMU.
    - ``prs.slides.add_slide(layout)`` adiciona um slide.
    - ``slide.shapes.add_picture(path, x, y, w, h)`` insere imagem.
    - ``slide.shapes.add_textbox(x, y, w, h)`` cria caixa de texto.
    - ``slide.shapes.add_group_shape()`` cria grupo (arrastar tudo junto).
"""

from __future__ import annotations

from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.shapes.shapetree import GroupShapes, SlideShapes
from pptx.util import Pt

from . import config as cfg
from .tipos import BlocoPessoa, LayoutMural, TextoPosicionado
from .utilitarios import px_fonte_para_pt, px_para_emu


# ───────────────────────────────────────────────────────────────────────────
# Calibragens específicas do PPTX
# ───────────────────────────────────────────────────────────────────────────
# O ``textbbox`` do Pillow devolve o bbox APERTADO dos glifos. O PowerPoint
# renderiza um pouco diferente (métricas variam entre engines), então a
# caixa de texto precisa de alguns pixels de folga para não cortar
# descendentes/ascendentes ou disparar quebra de linha indesejada.
# Estes paddings ficam isolados aqui para deixar claro que são ajuste fino.

PADDING_LARG_DIA_PX = 4.0     # folga ao redor do número do dia
PADDING_ALTURA_TEXTO_PX = 6.0  # folga vertical (linhas de dia/nome/cargo)
PADDING_TITULO_MES_PX = 16.0   # folga maior para a fonte grande do título

# Layout 6 da apresentação default do python-pptx = "Blank" (sem placeholders).
LAYOUT_BLANK_IDX = 6


# ───────────────────────────────────────────────────────────────────────────
# API pública
# ───────────────────────────────────────────────────────────────────────────


def exportar_pptx(
    layout: LayoutMural,
    caminho_template: Path,
    caminho_saida: Path,
) -> None:
    """Cria o ``.pptx`` editável a partir de um :class:`LayoutMural` pronto."""
    prs = Presentation()
    prs.slide_width = px_para_emu(layout.largura_imagem)
    prs.slide_height = px_para_emu(layout.altura_imagem)
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK_IDX])

    _adicionar_template_de_fundo(
        slide, caminho_template, prs.slide_width, prs.slide_height
    )
    _adicionar_titulo(slide.shapes, layout.titulo_mes)

    blocos_por_coluna: dict[int, list[BlocoPessoa]] = {}
    for bloco in layout.blocos:
        blocos_por_coluna.setdefault(bloco.coluna_idx, []).append(bloco)

    for idx_col in sorted(blocos_por_coluna):
        shapes_alvo = _abrir_destino_da_coluna(slide, idx_col)
        for bloco in blocos_por_coluna[idx_col]:
            for elemento in bloco.elementos:
                _adicionar_elemento(shapes_alvo, elemento)

    caminho_saida.parent.mkdir(parents=True, exist_ok=True)
    prs.save(caminho_saida)
    print(f"✅ Apresentação editável: {caminho_saida}")


# ───────────────────────────────────────────────────────────────────────────
# Internos — composição do slide
# ───────────────────────────────────────────────────────────────────────────


def _adicionar_template_de_fundo(slide, caminho_template: Path, largura_emu, altura_emu) -> None:
    """Adiciona o template como imagem cobrindo o slide todo."""
    slide.shapes.add_picture(
        str(caminho_template.resolve()),
        0,
        0,
        width=largura_emu,
        height=altura_emu,
    )


def _abrir_destino_da_coluna(slide, idx_col: int) -> Union[SlideShapes, GroupShapes]:
    """Devolve onde adicionar as caixas: o slide direto ou um grupo da coluna."""
    if not cfg.PPTX_GRUPO_POR_COLUNA:
        return slide.shapes
    grupo = slide.shapes.add_group_shape()
    grupo.name = f"Coluna {idx_col + 1}"
    return grupo.shapes


def _adicionar_titulo(shapes: SlideShapes, titulo: TextoPosicionado) -> None:
    """Caixa do título do mês — fora dos grupos das colunas."""
    _adicionar_caixa_texto(
        shapes,
        esquerda_px=float(titulo.x),
        topo_px=float(titulo.y),
        largura_px=float(titulo.largura) + PADDING_TITULO_MES_PX,
        altura_px=float(titulo.altura) + PADDING_TITULO_MES_PX,
        texto=titulo.texto,
        nome_fonte=cfg.PPTX_NOME_FONTE_MES,
        tamanho_pt=px_fonte_para_pt(cfg.TAMANHO_MES),
        cor=cfg.COR_TURQUESA,
        alinhamento=PP_ALIGN.CENTER,
    )


def _adicionar_elemento(
    shapes: Union[SlideShapes, GroupShapes],
    elemento: TextoPosicionado,
) -> None:
    """Cria a caixa de uma das partes (dia / nome / cargo)."""
    nome_fonte, tamanho_pt, cor = _estilo_para_papel(elemento.papel)

    # Para o dia mantemos uma caixa apertada em volta do número (só com a
    # folga de PADDING_LARG_DIA_PX). Para nome/cargo, garantimos que a
    # caixa caiba pelo menos a largura disponível na coluna — assim o
    # PowerPoint não quebra texto numa palavra que já cabia no PNG.
    if elemento.papel == "dia":
        largura = float(elemento.largura) + PADDING_LARG_DIA_PX
    else:
        largura = float(max(elemento.largura, elemento.largura_max_disponivel))

    _adicionar_caixa_texto(
        shapes,
        esquerda_px=float(elemento.x),
        topo_px=float(elemento.y),
        largura_px=largura,
        altura_px=float(elemento.altura) + PADDING_ALTURA_TEXTO_PX,
        texto=elemento.texto,
        nome_fonte=nome_fonte,
        tamanho_pt=tamanho_pt,
        cor=cor,
        alinhamento=PP_ALIGN.LEFT,
    )


def _estilo_para_papel(papel: str) -> tuple[str, float, tuple[int, int, int]]:
    if papel == "dia":
        return cfg.PPTX_NOME_FONTE_DIA, px_fonte_para_pt(cfg.TAMANHO_DIA), cfg.COR_TURQUESA
    if papel == "nome":
        return cfg.PPTX_NOME_FONTE_NOME, px_fonte_para_pt(cfg.TAMANHO_NOME), cfg.COR_PRETO
    if papel == "cargo":
        return cfg.PPTX_NOME_FONTE_CARGO, px_fonte_para_pt(cfg.TAMANHO_CARGO), cfg.COR_TURQUESA
    raise ValueError(f"Papel inesperado para PPTX: {papel!r}")


def _adicionar_caixa_texto(
    shapes: Union[SlideShapes, GroupShapes],
    *,
    esquerda_px: float,
    topo_px: float,
    largura_px: float,
    altura_px: float,
    texto: str,
    nome_fonte: str,
    tamanho_pt: float,
    cor: tuple[int, int, int],
    alinhamento: PP_ALIGN = PP_ALIGN.LEFT,
    negrito: bool = True,
) -> None:
    """Cria UMA caixa de texto com tipografia/posição já resolvidas.

    python-pptx em uso aqui:
      - ``shapes.add_textbox(left, top, width, height)`` recebe EMU.
      - ``caixa.text_frame`` é o "container" do texto; ``tf.clear()``
        zera o parágrafo default.
      - ``tf.paragraphs[0].font.{name,size,bold,color.rgb}`` define a
        tipografia. ``font.name`` precisa coincidir com a fonte instalada
        no Canva/PowerPoint para ser usada na hora de abrir o arquivo.
      - ``tf.word_wrap = True`` deixa o texto quebrar dentro da caixa.
      - ``margin_*`` zerados garantem que o pixel-position do PowerPoint
        bata com o calculado pelo Pillow (sem padding interno extra).
    """
    altura_px = max(altura_px, 1.0)
    largura_px = max(largura_px, 1.0)

    caixa = shapes.add_textbox(
        px_para_emu(esquerda_px),
        px_para_emu(topo_px),
        px_para_emu(largura_px),
        px_para_emu(altura_px),
    )
    tf = caixa.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.text = texto
    p.alignment = alinhamento
    if nome_fonte.strip():
        p.font.name = nome_fonte.strip()
    p.font.size = Pt(int(round(tamanho_pt)))
    p.font.bold = negrito
    p.font.color.rgb = RGBColor(cor[0], cor[1], cor[2])
